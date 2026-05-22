#!/usr/bin/env python3
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN

def create_pptx_template(path, title_text, subtitle_text):
    prs = Presentation()
    
    # Title Slide
    title_slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(title_slide_layout)
    title_shape = slide.shapes.title
    subtitle_shape = slide.placeholders[1]
    title_shape.text = title_text
    subtitle_shape.text = subtitle_text
    
    # Content Slide
    bullet_slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(bullet_slide_layout)
    title_shape = slide.shapes.title
    title_shape.text = "Project Key Highlights"
    
    body_shape = slide.placeholders[1]
    tf = body_shape.text_frame
    tf.text = "Point 1: Placeholder content here"
    p = tf.add_paragraph()
    p.text = "Point 2: Second key highlight"
    p = tf.add_paragraph()
    p.text = "Point 3: Third key highlight"
    
    # Closing Slide
    title_slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(title_slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    title.text = "Thank You"
    subtitle.text = "Questions & Discussion"
    
    prs.save(path)

if __name__ == "__main__":
    import sys
    if len(sys.argv) < 3:
        print("Usage: python3 generate_pptx.py <path> <title> <subtitle>")
        sys.exit(1)
    create_pptx_template(sys.argv[1], sys.argv[2], sys.argv[3])
