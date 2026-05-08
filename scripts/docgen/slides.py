#!/usr/bin/env python3
"""
slides.py — AInchors Presentation (PPTX)
Usage: python3 slides.py --title "Title" --output /path/out.pptx [--data /path/data.json]
"""
import argparse
import json
import sys
from datetime import date
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN


# AInchors brand colours
C_DARK   = RGBColor(0x0D, 0x11, 0x17)   # #0d1117
C_BLUE   = RGBColor(0x44, 0x72, 0xC4)   # #4472C4
C_WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
C_GREY   = RGBColor(0x88, 0x88, 0x88)
C_ACCENT = RGBColor(0x00, 0xB0, 0xF0)   # bright cyan accent

SLIDE_W = Inches(13.33)
SLIDE_H = Inches(7.5)


def add_rect(slide, left, top, width, height, color):
    from pptx.util import Emu
    shape = slide.shapes.add_shape(
        1,  # MSO_SHAPE_TYPE.RECTANGLE
        left, top, width, height
    )
    shape.line.fill.background()
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    return shape


def add_textbox(slide, text, left, top, width, height,
                font_size=18, bold=False, color=C_WHITE, align=PP_ALIGN.LEFT):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.color.rgb = color
    return txBox


def make_title_slide(prs, title, subtitle, doc_date):
    slide_layout = prs.slide_layouts[6]  # blank
    slide = prs.slides.add_slide(slide_layout)

    # Background
    add_rect(slide, 0, 0, SLIDE_W, SLIDE_H, C_DARK)

    # Accent line
    add_rect(slide, Inches(0.8), Inches(3.5), Inches(11.73), Inches(0.06), C_BLUE)

    # AInchors brand word
    add_textbox(slide, 'AInchors', Inches(0.8), Inches(0.4), Inches(5), Inches(0.8),
                font_size=22, bold=True, color=C_BLUE)

    # Title
    txBox = slide.shapes.add_textbox(Inches(0.8), Inches(1.4), Inches(11.73), Inches(1.8))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    run = p.add_run()
    run.text = title
    run.font.size = Pt(44)
    run.font.bold = True
    run.font.color.rgb = C_WHITE

    # Subtitle
    if subtitle:
        add_textbox(slide, subtitle, Inches(0.8), Inches(3.7), Inches(9), Inches(0.8),
                    font_size=20, color=C_GREY)

    # Date
    add_textbox(slide, doc_date, Inches(0.8), Inches(6.6), Inches(5), Inches(0.5),
                font_size=12, color=C_GREY)

    return slide


def make_content_slide(prs, slide_title, bullets):
    slide_layout = prs.slide_layouts[6]  # blank
    slide = prs.slides.add_slide(slide_layout)

    # White background
    add_rect(slide, 0, 0, SLIDE_W, SLIDE_H, C_WHITE)

    # Header bar
    add_rect(slide, 0, 0, SLIDE_W, Inches(1.1), C_DARK)

    # AInchors micro-brand in header
    add_textbox(slide, 'AInchors', Inches(0.3), Inches(0.15), Inches(2), Inches(0.5),
                font_size=11, bold=True, color=C_BLUE)

    # Slide title in header
    add_textbox(slide, slide_title, Inches(0.3), Inches(0.3), Inches(12.5), Inches(0.7),
                font_size=24, bold=True, color=C_WHITE)

    # Accent underline
    add_rect(slide, 0, Inches(1.1), SLIDE_W, Inches(0.04), C_BLUE)

    # Bullets
    if bullets:
        txBox = slide.shapes.add_textbox(Inches(0.6), Inches(1.35), Inches(12.1), Inches(5.8))
        tf = txBox.text_frame
        tf.word_wrap = True
        for i, bullet in enumerate(bullets):
            if i == 0:
                p = tf.paragraphs[0]
            else:
                p = tf.add_paragraph()
            p.alignment = PP_ALIGN.LEFT
            p.space_before = Pt(6)
            run = p.add_run()
            run.text = f'\u2022  {bullet}'
            run.font.size = Pt(18)
            run.font.color.rgb = C_DARK

    return slide


def main():
    parser = argparse.ArgumentParser()
    parser.add_argument('--title', required=True)
    parser.add_argument('--output', required=True)
    parser.add_argument('--data', default=None)
    args = parser.parse_args()

    data = {}
    if args.data:
        with open(args.data) as f:
            data = json.load(f)

    title = data.get('title', args.title)
    subtitle = data.get('subtitle', 'AInchors AI Transformation Consulting')
    doc_date = data.get('date', date.today().strftime('%d %B %Y'))
    slides_data = data.get('slides', [
        {
            'title': 'Engagement Overview',
            'bullets': [
                'AI readiness assessment across 5 business units',
                'Identified 12 high-impact automation opportunities',
                'Estimated annual savings: $2.4M+',
                'Timeline: 90-day discovery to recommendations',
            ]
        },
        {
            'title': 'Key Findings',
            'bullets': [
                'Manual processes account for 40% of operational overhead',
                'Data is siloed across 6 disconnected systems',
                'No current ML/AI capability in-house',
                'Strong leadership appetite for transformation',
            ]
        },
        {
            'title': 'Recommended Roadmap',
            'bullets': [
                'Phase 1 (Q3): Data unification & governance (8 weeks)',
                'Phase 2 (Q4): Process automation — top 3 use cases (12 weeks)',
                'Phase 3 (Q1 FY27): AI model deployment & upskilling (16 weeks)',
                'Ongoing: CoE establishment & continuous improvement',
            ]
        },
        {
            'title': 'Next Steps',
            'bullets': [
                'Sign engagement letter and confirm budget',
                'Schedule kick-off with executive sponsor',
                'Complete stakeholder interview plan (Week 1)',
                'Deliver Phase 1 roadmap document (Week 3)',
            ]
        },
    ])

    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    # Slide 1: Title
    make_title_slide(prs, title, subtitle, doc_date)

    # Content slides
    for s in slides_data:
        make_content_slide(prs, s.get('title', ''), s.get('bullets', []))

    prs.save(args.output)
    print(f'Slides saved: {args.output}')


if __name__ == '__main__':
    main()
